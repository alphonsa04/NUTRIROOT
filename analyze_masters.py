from pptx import Presentation
import os

def analyze_masters(path):
    prs = Presentation(path)
    print(f"Total Master slides: {len(prs.slide_masters)}")
    
    for i, master in enumerate(prs.slide_masters):
        print(f"\n--- Master {i} ---")
        for j, shape in enumerate(master.shapes):
            text = getattr(shape, "text", "No Text")
            print(f"Master Shape {j}: ID={shape.shape_id}, Name={shape.name}, Type={shape.shape_type}, Text='{text}'")
        
        for k, layout in enumerate(master.slide_layouts):
            print(f"  --- Layout {k}: {layout.name} ---")
            for l, shape in enumerate(layout.shapes):
                text = getattr(shape, "text", "No Text")
                if text != "No Text":
                    print(f"    Layout Shape {l}: ID={shape.shape_id}, Name={shape.name}, Text='{text}'")

if __name__ == "__main__":
    path = os.path.abspath("NutriRoot.pptx")
    analyze_masters(path)
