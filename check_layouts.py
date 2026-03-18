from pptx import Presentation
import os

def check_layout_shapes(path):
    prs = Presentation(path)
    print(f"Total Master slides: {len(prs.slide_masters)}")
    
    for i, master in enumerate(prs.slide_masters):
        print(f"\nMaster {i} Shapes:")
        for shape in master.shapes:
            print(f"  Shape: Name={shape.name}, Type={shape.shape_type}, Pos=({shape.left}, {shape.top})")
            
        for j, layout in enumerate(master.slide_layouts):
            print(f"Layout {j} '{layout.name}' Shapes:")
            for shape in layout.shapes:
                print(f"  Shape: Name={shape.name}, Type={shape.shape_type}, Pos=({shape.left}, {shape.top}), Size=({shape.width}, {shape.height})")

if __name__ == "__main__":
    path = os.path.abspath("NutriRoot.pptx")
    check_layout_shapes(path)
