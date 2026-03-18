from pptx import Presentation
import os

def find_gamma_deep(path):
    prs = Presentation(path)
    print(f"Total Slides: {len(prs.slides)}")
    
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            # Check text
            text = getattr(shape, "text", "").lower()
            if "gamma" in text:
                print(f"FOUND in Text on Slide {i+1}: '{getattr(shape, 'text', '')}' (Shape ID: {shape.shape_id})")
            
            # Check Alt Text
            alt_text = getattr(shape, "name", "").lower()
            if "gamma" in alt_text:
                print(f"FOUND in Name on Slide {i+1}: '{shape.name}' (Shape ID: {shape.shape_id})")
            
            # Check non-visible text (like in data structures if possible)
            try:
                if "gamma" in shape._element.xml.lower():
                    print(f"FOUND in XML on Slide {i+1}: Shape ID: {shape.shape_id}")
            except:
                pass

    # Also check Slide Masters and Layouts
    for i, master in enumerate(prs.slide_masters):
        for shape in master.shapes:
            if "gamma" in shape._element.xml.lower():
                print(f"FOUND in XML on Master {i}: Shape ID: {shape.shape_id}")
        for j, layout in enumerate(master.slide_layouts):
            for shape in layout.shapes:
                if "gamma" in shape._element.xml.lower():
                    print(f"FOUND in XML on Layout {i}-{j}: Shape ID: {shape.shape_id}")

if __name__ == "__main__":
    path = os.path.abspath("NutriRoot.pptx")
    find_gamma_deep(path)
