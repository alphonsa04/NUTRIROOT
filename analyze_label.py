from pptx import Presentation

def analyze_pptx(path):
    prs = Presentation(path)
    print(f"Total Slides: {len(prs.slides)}")
    
    # Check first slide thoroughly
    slide = prs.slides[0]
    print(f"\n--- Slide 1 ---")
    for i, shape in enumerate(slide.shapes):
        text = getattr(shape, "text", "No Text")
        print(f"Shape {i}: ID={shape.shape_id}, Name={shape.name}, Type={shape.shape_type}, Text='{text}'")

import os
if __name__ == "__main__":
    path = os.path.abspath("NutriRoot.pptx")
    analyze_pptx(path)
