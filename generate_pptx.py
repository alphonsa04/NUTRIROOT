from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "NutriRoot"
    subtitle.text = "Smart Farming Insights & Analytics"

    # Slide 2: Introduction
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Introduction"
    content.text = ("NutriRoot is a comprehensive web application designed to empower farmers with data-driven insights.\n"
                    "It addresses the challenge of optimizing crop yield by monitoring soil health and providing expert cultivation advice.")

    # Slide 3: Core Feature - Soil Monitoring
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Smart Soil Monitoring"
    content.text = ("Real-time tracking of essential soil parameters:\n"
                    "- NPK Levels (Nitrogen, Phosphorus, Potassium)\n"
                    "- Soil pH and Moisture content\n"
                    "- Ambient and Soil Temperature\n"
                    "- Automated alerts for critical conditions (e.g., Acidic Soil)")

    # Slide 4: Core Feature - Crop Knowledge
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Crop Knowledge & Recommendations"
    content.text = ("Expert guidance for optimized farming:\n"
                    "- Detailed tips for various crop species\n"
                    "- Optimal growth requirements (pH, nutrients, moisture)\n"
                    "- Market value insights (Demand: High/Medium/Low)\n"
                    "- Personalized crop recommendations based on soil data")

    # Slide 5: Core Feature - Agri-Marketplace
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Agri-Marketplace"
    content.text = ("Bridging the gap between farmers and quality supplies:\n"
                    "- Farmers can browse and purchase fertilizers and tools\n"
                    "- Sellers can list products and manage inventory\n"
                    "- Secure order processing and tracking\n"
                    "- Admin approval system for product quality control")

    # Slide 6: User Ecosystem
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "User Ecosystem"
    content.text = ("Tailored experiences for all stakeholders:\n"
                    "- Farmer Dashboard: Soil insights and marketplace access\n"
                    "- Seller Dashboard: Product management and order fulfilling\n"
                    "- Admin Panel: System analytics, user management, and approvals")

    # Slide 7: Technology Stack
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Technology Stack"
    content.text = ("Built with modern and robust technologies:\n"
                    "- Frontend: HTML5, CSS3, Vanilla JavaScript\n"
                    "- Backend & Database: Firebase (Auth, Firestore)\n"
                    "- IoT Integration: Python, MQTT Protocol\n"
                    "- Analytics: Python (Pandas, NumPy)")

    # Slide 8: System Architecture
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "System Architecture"
    content.text = ("1. IoT Sensors capture real-time soil data.\n"
                    "2. MQTT Bridge transmits data to the cloud.\n"
                    "3. Firebase stores and secures user and sensor data.\n"
                    "4. Web Application visualizes insights for farmers and admins.")

    # Slide 9: Innovation & Impact
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Innovation & Impact"
    content.text = ("- Reduces resource waste through precision farming\n"
                    "- Enhances food security by improving crop yields\n"
                    "- Simplifies the supply chain for agricultural products\n"
                    "- Digitizes traditional farming knowledge for accessibility")

    # Slide 10: Conclusion & Future Scope
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Conclusion & Future Scope"
    content.text = ("NutriRoot is the future of smart agriculture.\n"
                    "Future Scope:\n"
                    "- AI-powered pest and disease detection\n"
                    "- Integration with satellite imagery for large-scale analysis\n"
                    "- Expanded marketplace with logistics support")

    # Save the presentation
    output_path = "NutriRoot_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
