from pptx import Presentation
import os

def remove_branding_completely(path):
    prs = Presentation(path)
    removed_count = 0
    
    # Target position and size from analysis
    TARGET_LEFT = 12839215
    TARGET_TOP = 7749540
    TARGET_WIDTH = 1722605
    TARGET_HEIGHT = 411480
    
    # Tolerance for floating point comparison if needed, but here they are exact integers in the output
    
    def is_target(shape):
        return (shape.left == TARGET_LEFT and 
                shape.top == TARGET_TOP and 
                shape.width == TARGET_WIDTH and 
                shape.height == TARGET_HEIGHT)

    # 1. Remove from Slide Layouts
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            shapes_to_remove = [s for s in layout.shapes if is_target(s)]
            for shape in shapes_to_remove:
                try:
                    sp = shape._element
                    sp.getparent().remove(sp)
                    removed_count += 1
                    print(f"Removed branding from Layout: {layout.name}")
                except Exception as e:
                    print(f"Failed to remove from layout {layout.name}: {e}")
                    shape.left = -1000000 # Hide it

    # 2. Remove from Slides directly (if any)
    for i, slide in enumerate(prs.slides):
        shapes_to_remove = [s for s in slide.shapes if is_target(s)]
        for shape in shapes_to_remove:
            try:
                sp = shape._element
                sp.getparent().remove(sp)
                removed_count += 1
                print(f"Removed branding from Slide {i+1}")
            except Exception as e:
                print(f"Failed to remove from slide {i+1}: {e}")
                shape.left = -1000000

    if removed_count > 0:
        prs.save(path)
        print(f"Successfully removed {removed_count} instances of the branding label.")
    else:
        print("No branding labels found matching the target criteria.")

if __name__ == "__main__":
    path = os.path.abspath("NutriRoot.pptx")
    remove_branding_completely(path)
