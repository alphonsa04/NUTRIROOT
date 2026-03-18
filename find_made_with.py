from pptx import Presentation
import os

def find_made_with(path):
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            xml = shape._element.xml.lower()
            if "made with" in xml:
                print(f"FOUND 'made with' in XML on Slide {i+1}: Shape ID: {shape.shape_id}")
    
    for i, master in enumerate(prs.slide_masters):
        if "made with" in master._element.xml.lower():
            print(f"FOUND 'made with' in XML on Master {i}")
        for j, layout in enumerate(master.slide_layouts):
            if "made with" in layout._element.xml.lower():
                print(f"FOUND 'made with' in XML on Layout {i}-{j}")

if __name__ == "__main__":
    path = os.path.abspath("NutriRoot.pptx")
    find_made_with(path)
