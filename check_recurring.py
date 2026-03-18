from pptx import Presentation
import os

def check_recurring_shapes(path):
    prs = Presentation(path)
    all_slide_shapes = []
    
    print(f"Total Slides: {len(prs.slides)}")
    
    for i, slide in enumerate(prs.slides):
        print(f"\nSlide {i+1}:")
        for shape in slide.shapes:
            left, top = shape.left, shape.top
            width, height = shape.width, shape.height
            print(f"  Shape: ID={shape.shape_id}, Name={shape.name}, Type={shape.shape_type}, Pos=({left}, {top}), Size=({width}, {height})")

if __name__ == "__main__":
    path = os.path.abspath("NutriRoot.pptx")
    check_recurring_shapes(path)
