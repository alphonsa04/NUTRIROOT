from pptx import Presentation
import os

def find_common_shapes(path):
    prs = Presentation(path)
    if not prs.slides:
        return
    
    first_slide_shapes = []
    for shape in prs.slides[0].shapes:
        first_slide_shapes.append({
            'id': shape.shape_id,
            'name': shape.name,
            'pos': (shape.left, shape.top),
            'size': (shape.width, shape.height),
            'type': shape.shape_type,
            'count': 1
        })
    
    for i in range(1, len(prs.slides)):
        current_slide_shapes = prs.slides[i].shapes
        for f_shape in first_slide_shapes:
            for c_shape in current_slide_shapes:
                if (c_shape.left == f_shape['pos'][0] and 
                    c_shape.top == f_shape['pos'][1] and 
                    c_shape.width == f_shape['size'][0] and 
                    c_shape.height == f_shape['size'][1] and 
                    c_shape.shape_type == f_shape['type']):
                    f_shape['count'] += 1
                    break
    
    print(f"Total Slides: {len(prs.slides)}")
    for shape in first_slide_shapes:
        if shape['count'] >= len(prs.slides) - 1: # Common in almost all slides
            print(f"COMMON SHAPE: Name={shape['name']}, Pos={shape['pos']}, Size={shape['size']}, Count={shape['count']}")

if __name__ == "__main__":
    path = os.path.abspath("NutriRoot.pptx")
    find_common_shapes(path)
